from pptx import Presentation
import os
import glob
from copy import deepcopy  # 新增导入
from pptx.util import Pt  # 新增至顶部导入区域
from pptx.util import Inches

def copy_slide(target_prs, source_slide):
    """复制幻灯片及其内容到目标演示文稿"""
    target_slide = target_prs.slides.add_slide(source_slide.slide_layout)
    for shape in source_slide.shapes:
        el = shape.element
        new_el = deepcopy(el)  # 修改此处：使用deepcopy代替.copy()
        target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # 复制备注页（如有）
    if source_slide.notes_slide:
        notes = source_slide.notes_slide
        target_notes = target_slide.notes_slide
        for sh in notes.shapes:
            new_el = deepcopy(sh.element)  # 修改此处：使用deepcopy代替.copy()
            target_notes.shapes._spTree.insert_element_before(new_el, 'p:extLst')


def main():
    # 其他代码不变...
    # 定义输入输出路径
    input_dir = 'PPTS/'
    output_file = '20250414晨会点评_ABC.pptx'

    # 获取所有PPT文件列表
    ppt_files = glob.glob(os.path.join(input_dir, '*.ppt*'))
    if not ppt_files:
        print("警告：未找到任何PPT文件！")
        return

    # 创建新演示文稿
    new_prs = Presentation()

    # 创建第一页：文件列表展示
    first_layout = new_prs.slide_layouts[0]  # 使用标题幻灯片布局
    first_slide = new_prs.slides.add_slide(first_layout)
    title_shape = first_slide.shapes.title
    title_shape.top = Inches(3)
    title_shape.left = Inches(3)
    title_shape.text = "文件列表"

    text_box = first_slide.placeholders[1].text_frame
    text_box.clear()
    text_box.text = "\n".join([os.path.basename(f) for f in ppt_files])
    # 新增字体大小设置逻辑开始
    for paragraph in text_box.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)
    # 逐个复制所有PPT内容
    for file_path in ppt_files:
        src_prs = Presentation(file_path)
        for slide in src_prs.slides:
            copy_slide(new_prs, slide)

    # 保存最终结果
    new_prs.save(output_file)
    print(f"已完成合并，输出文件：{output_file}")

if __name__ == "__main__":
    main()

    # mmmu 任务